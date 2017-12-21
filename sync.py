from pptx import Presentation
from pprint import pprint
import json
import sys
import time, threading

def foo():
    print(time.ctime())
    sys.stdout.flush()

    data = json.load(open('data.json'))
    print(data)
    sys.stdout.flush()

    prs = Presentation('template.pptx')

    prs.slides[0].shapes.placeholders[0].text = data["title"]
    prs.slides[0].shapes.placeholders[1].text = data["content"]

    prs.save('test.pptx')

    threading.Timer(5, foo).start()

foo()

